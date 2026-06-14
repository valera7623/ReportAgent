#!/usr/bin/env python3
"""Create minimal PowerPoint template for ReportAgent."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def main() -> None:
    template_dir = Path(__file__).resolve().parent.parent / "app" / "templates"
    template_dir.mkdir(parents=True, exist_ok=True)
    out_path = template_dir / "template.pptx"

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "ReportAgent Report"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "%DATE%\nReportAgent"

    # Content slide with placeholders for metrics
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    content = prs.slides.add_slide(content_layout)
    if content.shapes.title:
        content.shapes.title.text = "Summary"
    body = content.placeholders[1] if len(content.placeholders) > 1 else None
    if body is not None:
        tf = body.text_frame
        tf.text = "%METRICS%"
        tf.paragraphs[0].font.size = Pt(14)

    # Blank slide for charts (Google Slides template uses text placeholders)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    chart_slide = prs.slides.add_slide(blank)
    box = chart_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    box.text_frame.text = "%CHART_1%"

    prs.save(out_path)
    print(f"Template saved: {out_path}")


if __name__ == "__main__":
    main()
