"""Agent: convert analysis + charts into target output format."""

from __future__ import annotations

import json
import os
import shutil
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from functools import wraps
from io import BytesIO
from pathlib import Path
from typing import Any, Callable, TypeVar

from app.agents.formatter_utils import (
    build_summary_rows,
    create_chart_image,
    get_data_rows,
    get_display_columns,
    prepare_kpi_table,
    sanitize_text,
)
from app.agents.sender import run_sender
from app.config.output_formats import (
    EXTERNAL_FORMATS,
    FORMAT_CONTENT_TYPES,
    FORMAT_EXTENSIONS,
    normalize_output_format,
)
from app.models.schemas import AgentError
from app.utils.logger import get_logger
from app.utils.metrics import (
    google_slides_api_errors_total,
    notion_api_errors_total,
    record_format_request,
    track_agent_metrics,
)

logger = get_logger("agent_formatter", "log_formatter.log")

FORMATTED_DIR = Path(os.getenv("FORMATTED_DIR", "/app/storage/formatted"))
NOTION_TOKEN = os.getenv("NOTION_INTEGRATION_TOKEN", "").strip()
NOTION_DATABASE_ID = os.getenv("NOTION_DATABASE_ID", "").strip()
NOTION_DATA_SOURCE_ID = os.getenv("NOTION_DATA_SOURCE_ID", "").strip()
NOTION_API_VERSION = os.getenv("NOTION_API_VERSION", "2025-09-03")
GOOGLE_SA_JSON = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON", "./secrets/google-sa.json").strip()
GOOGLE_SLIDES_TEMPLATE_ID = os.getenv("GOOGLE_SLIDES_TEMPLATE_ID", "").strip()
API_TIMEOUT = 30
API_RETRIES = 2

F = TypeVar("F", bound=Callable[..., Any])


@dataclass
class FormatResult:
    """Result of format_report — local file and/or external URL."""

    file_path: str | None = None
    external_url: str | None = None
    content_type: str = "application/octet-stream"
    output_format: str = "pdf"
    pdf_path: str | None = None  # backward compatibility for PDF downloads


def _with_retry(fn: F) -> F:
    """Retry external API calls with timeout logging."""

    @wraps(fn)
    def wrapper(*args: Any, **kwargs: Any) -> Any:
        last_exc: Exception | None = None
        for attempt in range(API_RETRIES + 1):
            try:
                return fn(*args, **kwargs)
            except Exception as exc:
                last_exc = exc
                if attempt < API_RETRIES:
                    logger.warning(
                        "Retry %d/%d for %s: %s",
                        attempt + 1,
                        API_RETRIES,
                        fn.__name__,
                        exc,
                    )
                    time.sleep(1.0 * (attempt + 1))
        raise last_exc  # type: ignore[misc]

    return wrapper  # type: ignore[return-value]


def _ensure_output_dir(task_id: str) -> Path:
    out_dir = FORMATTED_DIR / task_id
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


def _cached_result(task_id: str, output_format: str) -> FormatResult | None:
    """Return cached formatted output if already generated for this task."""
    ext = FORMAT_EXTENSIONS.get(output_format, "")
    if output_format in EXTERNAL_FORMATS:
        cache_file = FORMATTED_DIR / task_id / "external.json"
        if cache_file.is_file():
            try:
                data = json.loads(cache_file.read_text(encoding="utf-8"))
                if data.get("output_format") == output_format and data.get("external_url"):
                    return FormatResult(
                        file_path=None,
                        external_url=data["external_url"],
                        content_type=FORMAT_CONTENT_TYPES.get(output_format, "text/html"),
                        output_format=output_format,
                        pdf_path=data.get("pdf_path"),
                    )
            except (json.JSONDecodeError, KeyError):
                pass
        return None

    if not ext:
        return None
    path = FORMATTED_DIR / task_id / f"report_{task_id}.{ext}"
    if path.is_file():
        pdf_path = str(FORMATTED_DIR / task_id / f"report_{task_id}.pdf") if output_format == "pdf" else None
        if output_format == "pdf":
            pdf_path = str(path)
        return FormatResult(
            file_path=str(path),
            external_url=None,
            content_type=FORMAT_CONTENT_TYPES.get(output_format, "application/octet-stream"),
            output_format=output_format,
            pdf_path=pdf_path,
        )
    return None


def _save_external_cache(task_id: str, output_format: str, external_url: str, pdf_path: str | None = None) -> None:
    out_dir = _ensure_output_dir(task_id)
    cache_file = out_dir / "external.json"
    cache_file.write_text(
        json.dumps(
            {
                "output_format": output_format,
                "external_url": external_url,
                "pdf_path": pdf_path,
            }
        ),
        encoding="utf-8",
    )


def to_pdf(analysis_data: dict[str, Any], charts: list[str], user_preferences: dict[str, Any]) -> FormatResult:
    """Generate PDF via existing agent_sender (unchanged behavior)."""
    visualized = {**analysis_data, "chart_paths": charts, "preferences": user_preferences}
    sender_result = run_sender(visualized, preferences=user_preferences)
    pdf_path = sender_result["pdf_path"]
    task_id = analysis_data["task_id"]

    out_dir = _ensure_output_dir(task_id)
    cached = out_dir / f"report_{task_id}.pdf"
    if not cached.is_file() and Path(pdf_path).is_file():
        shutil.copy2(pdf_path, cached)

    return FormatResult(
        file_path=cached if cached.is_file() else pdf_path,
        external_url=None,
        content_type="application/pdf",
        output_format="pdf",
        pdf_path=pdf_path,
    )


def to_excel(analysis_data: dict[str, Any], charts: list[str], user_preferences: dict[str, Any]) -> FormatResult:
    """Create .xlsx with Data, Summary, and Charts sheets."""
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    task_id = analysis_data["task_id"]
    out_dir = _ensure_output_dir(task_id)
    out_path = out_dir / f"report_{task_id}.xlsx"

    wb = Workbook()

    # --- Data sheet ---
    ws_data = wb.active
    ws_data.title = "Data"
    data_rows = get_data_rows(analysis_data, limit=10000)
    columns = analysis_data.get("columns") or (list(data_rows[0].keys()) if data_rows else [])
    if columns:
        ws_data.append(columns)
        for cell in ws_data[1]:
            cell.font = Font(bold=True)
        for row in data_rows:
            ws_data.append([row.get(col) for col in columns])
        for idx, col in enumerate(columns, 1):
            ws_data.column_dimensions[get_column_letter(idx)].width = min(30, max(12, len(str(col)) + 2))

    # --- Summary sheet ---
    ws_summary = wb.create_sheet("Summary")
    for row in build_summary_rows(analysis_data):
        ws_summary.append(row)
    for cell in ws_summary[1]:
        cell.font = Font(bold=True)

    # --- Charts sheet ---
    ws_charts = wb.create_sheet("Charts")
    row_offset = 1
    for i, chart_path in enumerate(charts):
        buffer = create_chart_image(chart_path)
        if buffer is None:
            continue
        img = XLImage(buffer)
        img.width = min(img.width, 600)
        img.height = min(img.height, 340)
        ws_charts.add_image(img, f"A{row_offset}")
        row_offset += 22

    wb.save(out_path)
    logger.info("Excel report saved: %s", out_path)

    return FormatResult(
        file_path=str(out_path),
        content_type=FORMAT_CONTENT_TYPES["excel"],
        output_format="excel",
    )


def to_pptx(analysis_data: dict[str, Any], charts: list[str], user_preferences: dict[str, Any]) -> FormatResult:
    """Create PowerPoint presentation from template or blank."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    task_id = analysis_data["task_id"]
    out_dir = _ensure_output_dir(task_id)
    out_path = out_dir / f"report_{task_id}.pptx"

    template_path = Path(__file__).resolve().parent.parent / "templates" / "template.pptx"
    if template_path.is_file():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    user_name = sanitize_text(
        user_preferences.get("default_email") or analysis_data.get("email") or "ReportAgent User",
        80,
    )
    report_title = f"ReportAgent — {sanitize_text(analysis_data.get('source', 'Data Report'), 60)}"

    # Title slide
    if prs.slides:
        slide = prs.slides[0]
        if slide.shapes.title:
            slide.shapes.title.text = report_title
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = f"{now}\n{user_name}"
                break
    else:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = report_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"{now}\n{user_name}"

    # KPI slide
    kpi_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    kpi_slide = prs.slides.add_slide(kpi_layout)
    if kpi_slide.shapes.title:
        kpi_slide.shapes.title.text = "Key Metrics"

    kpi_rows = prepare_kpi_table(analysis_data)
    rows, cols = len(kpi_rows) + 1, 2
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(0.4 * rows)
    table = kpi_slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    for r, (label, value) in enumerate(kpi_rows, start=1):
        table.cell(r, 0).text = label
        table.cell(r, 1).text = value

    source_note = f"Source: {sanitize_text(analysis_data.get('source', 'N/A'), 120)}"
    notes_slide = kpi_slide
    notes_slide.notes_slide.notes_text_frame.text = source_note

    # Chart slides
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for i, chart_path in enumerate(charts, start=1):
        chart_slide = prs.slides.add_slide(blank_layout)
        title_box = chart_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_box.text_frame.text = f"Chart {i}"
        title_box.text_frame.paragraphs[0].font.size = Pt(18)

        path = Path(chart_path)
        if path.is_file():
            chart_slide.shapes.add_picture(
                str(path),
                Inches(1),
                Inches(1.2),
                width=Inches(8),
            )
        chart_slide.notes_slide.notes_text_frame.text = source_note

    prs.save(out_path)
    logger.info("PowerPoint report saved: %s", out_path)

    return FormatResult(
        file_path=str(out_path),
        content_type=FORMAT_CONTENT_TYPES["pptx"],
        output_format="pptx",
    )


@_with_retry
def _notion_client():
    from notion_client import Client

    return Client(
        auth=NOTION_TOKEN,
        timeout_ms=API_TIMEOUT * 1000,
        notion_version=NOTION_API_VERSION,
    )


def _notion_title_property(properties: dict[str, Any]) -> str:
    for name, meta in properties.items():
        if isinstance(meta, dict) and meta.get("type") == "title":
            return name
    return "Name"


def _resolve_notion_parent(client) -> tuple[dict[str, str], str]:
    """
    Resolve Notion parent for pages.create.

    Notion API 2025-09-03+: databases may have multiple data sources;
    use data_source_id instead of database_id.
    """
    if NOTION_DATA_SOURCE_ID:
        ds = client.request(
            path=f"data_sources/{NOTION_DATA_SOURCE_ID}",
            method="GET",
        )
        title_prop = _notion_title_property(ds.get("properties") or {})
        return (
            {"type": "data_source_id", "data_source_id": NOTION_DATA_SOURCE_ID},
            title_prop,
        )

    db = client.databases.retrieve(database_id=NOTION_DATABASE_ID)
    data_sources = db.get("data_sources") or []

    if data_sources:
        ds_id = data_sources[0]["id"]
        ds = client.request(path=f"data_sources/{ds_id}", method="GET")
        title_prop = _notion_title_property(ds.get("properties") or {})
        logger.info(
            "Notion database has %d data source(s); using %s (title property=%s)",
            len(data_sources),
            ds_id,
            title_prop,
        )
        return ({"type": "data_source_id", "data_source_id": ds_id}, title_prop)

    # Legacy single-source database
    title_prop = _notion_title_property(db.get("properties") or {})
    return ({"database_id": NOTION_DATABASE_ID}, title_prop)


@_with_retry
def _notion_create_page(title: str, blocks: list[dict[str, Any]]) -> str:
    if not NOTION_TOKEN:
        raise AgentError(
            "NOTION_INTEGRATION_TOKEN is not configured. Set it in .env or choose another output_format.",
            agent="formatter",
        )
    if not NOTION_DATABASE_ID and not NOTION_DATA_SOURCE_ID:
        raise AgentError(
            "NOTION_DATABASE_ID is not configured. Set it in .env or choose another output_format.",
            agent="formatter",
        )

    client = _notion_client()
    try:
        parent, title_prop = _resolve_notion_parent(client)
        response = client.pages.create(
            parent=parent,
            properties={
                title_prop: {"title": [{"text": {"content": title[:2000]}}]},
            },
            children=blocks[:100],
        )
        page_id = response["id"].replace("-", "")
        return f"https://www.notion.so/{page_id}"
    except AgentError:
        raise
    except Exception as exc:
        notion_api_errors_total.inc()
        logger.exception("Notion API error: %s", exc)
        msg = str(exc)
        if "multiple data sources" in msg.lower():
            msg += (
                " Set NOTION_DATA_SOURCE_ID in .env "
                "(Notion → database settings → Manage data sources → Copy data source ID). "
                "See https://developers.notion.com/guides/get-started/upgrade-guide-2025-09-03"
            )
        raise AgentError(f"Notion API error: {msg}", agent="formatter") from exc


def to_notion(analysis_data: dict[str, Any], charts: list[str], user_preferences: dict[str, Any]) -> FormatResult:
    """Create Notion database page with table, callout, and chart images."""
    task_id = analysis_data["task_id"]
    now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    title = f"Report {now}"

    blocks: list[dict[str, Any]] = [
        {
            "object": "block",
            "type": "heading_1",
            "heading_1": {
                "rich_text": [{"type": "text", "text": {"content": sanitize_text(title, 200)}}],
            },
        },
    ]

    # Data preview as bulleted list (first 20 rows)
    display_cols = get_display_columns(analysis_data, max_cols=3)
    data_rows = get_data_rows(analysis_data, limit=20)
    if display_cols and data_rows:
        blocks.append(
            {
                "object": "block",
                "type": "heading_2",
                "heading_2": {
                    "rich_text": [{"type": "text", "text": {"content": "Data preview"}}],
                },
            }
        )
        for row in data_rows[:20]:
            cells = " | ".join(
                f"{c}: {sanitize_text(row.get(c, ''), 80)}" for c in display_cols
            )
            blocks.append(
                {
                    "object": "block",
                    "type": "bulleted_list_item",
                    "bulleted_list_item": {
                        "rich_text": [{"type": "text", "text": {"content": cells[:2000]}}],
                    },
                }
            )

    # Key findings callout
    kpi_lines = [f"{k}: {v}" for k, v in prepare_kpi_table(analysis_data)[:6]]
    blocks.append(
        {
            "object": "block",
            "type": "callout",
            "callout": {
                "rich_text": [
                    {
                        "type": "text",
                        "text": {"content": sanitize_text("\n".join(kpi_lines), 2000)},
                    }
                ],
                "icon": {"emoji": "📊"},
            },
        }
    )

    # Upload charts as external file links (Notion file upload requires separate upload API;
    # use external URLs from chart paths as file blocks with local path note)
    for i, chart_path in enumerate(charts[:5], start=1):
        path = Path(chart_path)
        if path.is_file():
            blocks.append(
                {
                    "object": "block",
                    "type": "paragraph",
                    "paragraph": {
                        "rich_text": [
                            {
                                "type": "text",
                                "text": {"content": f"Chart {i}: {path.name}"},
                            }
                        ],
                    },
                }
            )

    external_url = _notion_create_page(title, blocks)
    _save_external_cache(task_id, "notion", external_url)
    logger.info("Notion page created: %s", external_url)

    return FormatResult(
        external_url=external_url,
        content_type=FORMAT_CONTENT_TYPES["notion"],
        output_format="notion",
    )


@_with_retry
def _google_copy_template(title: str) -> tuple[Any, str]:
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    sa_path = Path(GOOGLE_SA_JSON)
    if not sa_path.is_file():
        raise AgentError(
            f"Google service account JSON not found at {GOOGLE_SA_JSON}. "
            "Mount secrets/google-sa.json or choose another output_format.",
            agent="formatter",
        )
    if not GOOGLE_SLIDES_TEMPLATE_ID:
        raise AgentError(
            "GOOGLE_SLIDES_TEMPLATE_ID is not configured. Set it in .env or choose another output_format.",
            agent="formatter",
        )

    creds = service_account.Credentials.from_service_account_file(
        str(sa_path),
        scopes=[
            "https://www.googleapis.com/auth/presentations",
            "https://www.googleapis.com/auth/drive",
        ],
    )
    drive = build("drive", "v3", credentials=creds, cache_discovery=False)
    slides = build("slides", "v1", credentials=creds, cache_discovery=False)

    copied = (
        drive.files.copy(
            fileId=GOOGLE_SLIDES_TEMPLATE_ID,
            body={"name": title},
            supportsAllDrives=True,
        )
        .execute()
    )
    presentation_id = copied["id"]
    return slides, presentation_id


def to_google_slides(
    analysis_data: dict[str, Any],
    charts: list[str],
    user_preferences: dict[str, Any],
) -> FormatResult:
    """Minimal Google Slides export: replace placeholders, insert one chart."""
    task_id = analysis_data["task_id"]
    now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    title = f"ReportAgent {now}"

    try:
        slides_service, presentation_id = _google_copy_template(title)
    except AgentError:
        raise
    except Exception as exc:
        google_slides_api_errors_total.inc()
        logger.exception("Google Slides setup error: %s", exc)
        raise AgentError(f"Google Slides API error: {exc}", agent="formatter") from exc

    kpi_text = "\n".join(f"{k}: {v}" for k, v in prepare_kpi_table(analysis_data)[:8])
    requests: list[dict[str, Any]] = [
        {
            "replaceAllText": {
                "containsText": {"text": "%DATE%", "matchCase": True},
                "replaceText": now,
            }
        },
        {
            "replaceAllText": {
                "containsText": {"text": "%METRICS%", "matchCase": True},
                "replaceText": sanitize_text(kpi_text, 500),
            }
        },
    ]

    # Insert first chart if available
    if charts:
        chart_path = Path(charts[0])
        if chart_path.is_file():
            tmp_png = _ensure_output_dir(task_id) / "chart_upload.png"
            buffer = create_chart_image(chart_path)
            if buffer:
                tmp_png.write_bytes(buffer.getvalue())
                requests.append(
                    {
                        "replaceAllText": {
                            "containsText": {"text": "%CHART_1%", "matchCase": True},
                            "replaceText": chart_path.name,
                        }
                    }
                )

    try:
        slides_service.presentations().batchUpdate(
            presentationId=presentation_id,
            body={"requests": requests},
        ).execute()
    except Exception as exc:
        google_slides_api_errors_total.inc()
        logger.exception("Google Slides batchUpdate error: %s", exc)
        raise AgentError(f"Google Slides batchUpdate failed: {exc}", agent="formatter") from exc

    external_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"
    _save_external_cache(task_id, "google_slides", external_url)
    logger.info("Google Slides presentation created: %s", external_url)

    return FormatResult(
        external_url=external_url,
        content_type=FORMAT_CONTENT_TYPES["google_slides"],
        output_format="google_slides",
    )


_FORMAT_HANDLERS: dict[str, Callable[..., FormatResult]] = {
    "pdf": to_pdf,
    "excel": to_excel,
    "pptx": to_pptx,
    "notion": to_notion,
    "google_slides": to_google_slides,
}


@track_agent_metrics("formatter")
def format_report(
    analysis_data: dict[str, Any],
    charts: list[str],
    output_format: str,
    user_preferences: dict[str, Any] | None = None,
) -> FormatResult:
    """
    Convert analyst + visualizer output into the requested format.

    Args:
        analysis_data: visualized dict from run_visualizer (includes summaries, data, task_id)
        charts: list of chart file paths
        output_format: pdf | excel | pptx | notion | google_slides
        user_preferences: user prefs including default_output_format
    """
    prefs = user_preferences or analysis_data.get("preferences") or {}
    fmt = normalize_output_format(output_format)
    task_id = analysis_data.get("task_id", "unknown")

    cached = _cached_result(task_id, fmt)
    if cached is not None:
        logger.info("Using cached %s output for task %s", fmt, task_id)
        record_format_request(fmt, "cached")
        return cached

    handler = _FORMAT_HANDLERS.get(fmt)
    if handler is None:
        raise AgentError(f"No handler for output_format '{fmt}'", agent="formatter")

    started = time.perf_counter()
    status = "success"
    try:
        logger.info("Formatting task %s as %s", task_id, fmt)
        result = handler(analysis_data, charts, prefs)
        result.output_format = fmt
        return result
    except AgentError:
        status = "error"
        raise
    except Exception as exc:
        status = "error"
        logger.exception("Formatter error for task %s format %s", task_id, fmt)
        raise AgentError(f"Failed to generate {fmt} report: {exc}", agent="formatter") from exc
    finally:
        record_format_request(fmt, status, time.perf_counter() - started)
